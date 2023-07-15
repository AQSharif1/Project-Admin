import requests
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
import pandas as pd
from sqlalchemy import create_engine
import os
import pptx.util
from collections.abc import Container

# Make a GET request to the API
response = requests.get("https://jsonplaceholder.typicode.com/todos")

# Check if the request was successful
if response.status_code == 200:
    # Get the JSON data from the response
    data = response.json()

    # Create a DataFrame from the JSON data
    df = pd.DataFrame(data)

    # Display the DataFrame
    print(df)
else:
    print("Error occurred while accessing the API:", response.status_code)

def buildPresentation(df):
   API_BASE = "https://abcd2.projectabcd.com/api/getinfo.php?id="
   print("Creating powerpoint slides.")
   
   with open("preferences.txt", "r") as f:
        slideOption = 1
        textFont = f.readline().split("= ")
        textFont = textFont[1]
        titleFont = f.readline().split("= ")
        titleFont = titleFont[1]
        textSize = f.readline().split("= ")
        textSize = int(textSize[1])
        titleSize = f.readline().split("= ")
        titleSize = int(titleSize[1])
        prs = Presentation()
        presentationLength = 10
        

   for i in range(0,presentationLength): 
    (df['id'][i]) 
    API_BASE = API_BASE + str(df['id'][i])   
    print("value: " + str(API_BASE))  
    prs.slide_width = pptx.util.Inches(8)
    prs.slide_height = pptx.util.Inches(11)  
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    API_BASE_slide = slide2.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width=pptx.util.Inches(6),height=pptx.util.Inches(7))   
    contentBoxtf = API_BASE_slide.text_frame
    contentBoxtf.word_wrap = True
   test = "test_excel.pptx"
   prs.save(test)
   return test

test = buildPresentation(df)
os.startfile(test)
